import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import os
import magic
from PIL import Image, ImageTk
import subprocess
from pathlib import Path
import threading
import json
import pandas as pd
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
import pypandoc
from io import BytesIO

class FileConverter:
    def __init__(self, root):
        self.root = root
        self.root.title("文件格式转换器")
        self.root.geometry("800x600")
        self.root.resizable(True, True)
        
        # 当前文件信息
        self.current_file = None
        self.current_format = None
        
        # 支持的转换格式映射
        self.format_mappings = {
            'image': ['PNG', 'JPEG', 'GIF', 'BMP', 'TIFF', 'WEBP', 'ICO'],
            'document': ['PDF', 'DOCX', 'DOC', 'TXT', 'RTF', 'HTML', 'MD'],
            'spreadsheet': ['XLSX', 'XLS', 'CSV', 'ODS', 'TSV'],
            'presentation': ['PPTX', 'PPT', 'ODP', 'PDF'],
            'audio': ['MP3', 'WAV', 'FLAC', 'AAC', 'OGG'],
            'video': ['MP4', 'AVI', 'MOV', 'MKV', 'WEBM', 'FLV'],
            'archive': ['ZIP', 'RAR', '7Z', 'TAR.GZ']
        }
        
        self.setup_ui()
        
    def setup_ui(self):
        """设置用户界面"""
        # 主框架
        main_frame = ttk.Frame(self.root, padding="10")
        main_frame.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        
        # 配置网格权重
        self.root.columnconfigure(0, weight=1)
        self.root.rowconfigure(0, weight=1)
        main_frame.columnconfigure(1, weight=1)
        
        # 文件选择区域
        file_frame = ttk.LabelFrame(main_frame, text="文件选择", padding="10")
        file_frame.grid(row=0, column=0, columnspan=2, sticky=(tk.W, tk.E), pady=(0, 10))
        file_frame.columnconfigure(1, weight=1)
        
        ttk.Button(file_frame, text="选择文件", command=self.select_file).grid(row=0, column=0, padx=(0, 10))
        
        self.file_path_var = tk.StringVar()
        self.file_path_entry = ttk.Entry(file_frame, textvariable=self.file_path_var, state="readonly")
        self.file_path_entry.grid(row=0, column=1, sticky=(tk.W, tk.E))
        
        # 文件信息区域
        info_frame = ttk.LabelFrame(main_frame, text="文件信息", padding="10")
        info_frame.grid(row=1, column=0, columnspan=2, sticky=(tk.W, tk.E), pady=(0, 10))
        info_frame.columnconfigure(1, weight=1)
        
        ttk.Label(info_frame, text="检测到的格式:").grid(row=0, column=0, sticky=tk.W, padx=(0, 10))
        self.detected_format_var = tk.StringVar(value="未选择文件")
        ttk.Label(info_frame, textvariable=self.detected_format_var, font=("Arial", 10, "bold")).grid(row=0, column=1, sticky=tk.W)
        
        ttk.Label(info_frame, text="文件大小:").grid(row=1, column=0, sticky=tk.W, padx=(0, 10))
        self.file_size_var = tk.StringVar(value="-")
        ttk.Label(info_frame, textvariable=self.file_size_var).grid(row=1, column=1, sticky=tk.W)
        
        # 转换设置区域
        convert_frame = ttk.LabelFrame(main_frame, text="转换设置", padding="10")
        convert_frame.grid(row=2, column=0, columnspan=2, sticky=(tk.W, tk.E), pady=(0, 10))
        convert_frame.columnconfigure(1, weight=1)
        
        ttk.Label(convert_frame, text="目标格式:").grid(row=0, column=0, sticky=tk.W, padx=(0, 10))
        self.target_format_var = tk.StringVar()
        self.target_format_combo = ttk.Combobox(convert_frame, textvariable=self.target_format_var, state="readonly")
        self.target_format_combo.grid(row=0, column=1, sticky=(tk.W, tk.E), padx=(0, 10))
        
        ttk.Label(convert_frame, text="输出目录:").grid(row=1, column=0, sticky=tk.W, padx=(0, 10), pady=(10, 0))
        
        output_frame = ttk.Frame(convert_frame)
        output_frame.grid(row=1, column=1, sticky=(tk.W, tk.E), pady=(10, 0))
        output_frame.columnconfigure(0, weight=1)
        
        self.output_path_var = tk.StringVar()
        self.output_path_entry = ttk.Entry(output_frame, textvariable=self.output_path_var)
        self.output_path_entry.grid(row=0, column=0, sticky=(tk.W, tk.E), padx=(0, 10))
        
        ttk.Button(output_frame, text="浏览", command=self.select_output_dir).grid(row=0, column=1)
        
        # 转换按钮
        self.convert_button = ttk.Button(convert_frame, text="开始转换", command=self.start_conversion, state="disabled")
        self.convert_button.grid(row=2, column=0, columnspan=2, pady=(20, 0))
        
        # 进度条
        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(main_frame, variable=self.progress_var, maximum=100)
        self.progress_bar.grid(row=3, column=0, columnspan=2, sticky=(tk.W, tk.E), pady=(0, 10))
        
        # 状态栏
        self.status_var = tk.StringVar(value="准备就绪")
        status_label = ttk.Label(main_frame, textvariable=self.status_var, relief=tk.SUNKEN, anchor=tk.W)
        status_label.grid(row=4, column=0, columnspan=2, sticky=(tk.W, tk.E))
        
    def select_file(self):
        """选择要转换的文件"""
        file_path = filedialog.askopenfilename(
            title="选择要转换的文件",
            filetypes=[
                ("所有文件", "*.*"),
                ("图片文件", "*.png *.jpg *.jpeg *.gif *.bmp *.tiff *.webp"),
                ("文档文件", "*.pdf *.docx *.doc *.txt *.rtf *.html *.md"),
                ("电子表格", "*.xlsx *.xls *.csv *.ods *.tsv"),
                ("演示文稿", "*.pptx *.ppt *.odp"),
                ("音频文件", "*.mp3 *.wav *.flac *.aac *.ogg"),
                ("视频文件", "*.mp4 *.avi *.mov *.mkv *.webm *.flv"),
            ]
        )
        
        if file_path:
            self.current_file = file_path
            self.file_path_var.set(file_path)
            self.detect_file_format()
            
    def detect_file_format(self):
        """检测文件格式"""
        if not self.current_file:
            return
            
        try:
            # 获取文件信息
            file_stat = os.stat(self.current_file)
            file_size = self.format_file_size(file_stat.st_size)
            self.file_size_var.set(file_size)
            
            # 使用python-magic检测文件类型
            try:
                mime_type = magic.from_file(self.current_file, mime=True)
                file_type = magic.from_file(self.current_file)
            except:
                # 如果magic失败，使用文件扩展名
                ext = os.path.splitext(self.current_file)[1].lower()
                mime_type = self.guess_mime_type(ext)
                file_type = f"文件扩展名: {ext}"
            
            self.current_format = mime_type
            self.detected_format_var.set(f"{file_type}")
            
            # 更新可用的转换格式
            self.update_target_formats(mime_type)
            
            # 设置默认输出目录
            if not self.output_path_var.get():
                output_dir = os.path.dirname(self.current_file)
                self.output_path_var.set(output_dir)
                
            self.status_var.set("文件检测完成，请选择目标格式")
            
        except Exception as e:
            messagebox.showerror("错误", f"文件检测失败: {str(e)}")
            
    def guess_mime_type(self, ext):
        """根据文件扩展名猜测MIME类型"""
        mime_map = {
            '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg', '.png': 'image/png',
            '.gif': 'image/gif', '.bmp': 'image/bmp', '.tiff': 'image/tiff',
            '.webp': 'image/webp', '.ico': 'image/x-icon',
            '.pdf': 'application/pdf', 
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.doc': 'application/msword',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.xls': 'application/vnd.ms-excel',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.ppt': 'application/vnd.ms-powerpoint',
            '.csv': 'text/csv', '.tsv': 'text/tab-separated-values',
            '.ods': 'application/vnd.oasis.opendocument.spreadsheet',
            '.odp': 'application/vnd.oasis.opendocument.presentation',
            '.txt': 'text/plain', '.rtf': 'application/rtf', '.html': 'text/html',
            '.md': 'text/markdown',
            '.mp3': 'audio/mpeg', '.wav': 'audio/wav', '.flac': 'audio/flac',
            '.aac': 'audio/aac', '.ogg': 'audio/ogg',
            '.mp4': 'video/mp4', '.avi': 'video/x-msvideo', '.mov': 'video/quicktime',
            '.mkv': 'video/x-matroska', '.webm': 'video/webm', '.flv': 'video/x-flv',
            '.zip': 'application/zip', '.rar': 'application/x-rar-compressed',
            '.7z': 'application/x-7z-compressed'
        }
        return mime_map.get(ext, 'application/octet-stream')
        
    def update_target_formats(self, mime_type):
        """根据检测到的文件类型更新可用的转换格式"""
        formats = []
        
        if mime_type.startswith('image/'):
            formats = self.format_mappings['image']
        elif mime_type.startswith('audio/'):
            formats = self.format_mappings['audio']
        elif mime_type.startswith('video/'):
            formats = self.format_mappings['video']
        elif mime_type in ['application/pdf', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', 
                          'application/msword', 'text/plain', 'application/rtf', 'text/html', 'text/markdown']:
            formats = self.format_mappings['document']
        elif mime_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 
                          'application/vnd.ms-excel', 'text/csv', 'text/tab-separated-values',
                          'application/vnd.oasis.opendocument.spreadsheet']:
            formats = self.format_mappings['spreadsheet']
        elif mime_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation',
                          'application/vnd.ms-powerpoint', 'application/vnd.oasis.opendocument.presentation']:
            formats = self.format_mappings['presentation']
        elif mime_type in ['application/zip', 'application/x-rar-compressed', 'application/x-7z-compressed']:
            formats = self.format_mappings['archive']
        else:
            # 默认提供一些通用格式
            formats = ['TXT', 'PDF']
            
        self.target_format_combo['values'] = formats
        if formats:
            self.target_format_combo.set(formats[0])
            self.convert_button['state'] = 'normal'
        else:
            self.convert_button['state'] = 'disabled'
            
    def select_output_dir(self):
        """选择输出目录"""
        output_dir = filedialog.askdirectory(title="选择输出目录")
        if output_dir:
            self.output_path_var.set(output_dir)
            
    def format_file_size(self, size_bytes):
        """格式化文件大小"""
        if size_bytes == 0:
            return "0 B"
        
        size_names = ["B", "KB", "MB", "GB", "TB"]
        import math
        i = int(math.floor(math.log(size_bytes, 1024)))
        p = math.pow(1024, i)
        s = round(size_bytes / p, 2)
        return f"{s} {size_names[i]}"
        
    def start_conversion(self):
        """开始转换过程"""
        if not self.current_file or not self.target_format_var.get() or not self.output_path_var.get():
            messagebox.showerror("错误", "请确保已选择文件、目标格式和输出目录")
            return
            
        # 在新线程中执行转换以避免界面冻结
        self.convert_button['state'] = 'disabled'
        self.progress_var.set(0)
        
        conversion_thread = threading.Thread(target=self.perform_conversion)
        conversion_thread.daemon = True
        conversion_thread.start()
        
    def perform_conversion(self):
        """执行实际的文件转换"""
        try:
            source_file = self.current_file
            target_format = self.target_format_var.get().lower()
            output_dir = self.output_path_var.get()
            
            # 生成输出文件名
            base_name = os.path.splitext(os.path.basename(source_file))[0]
            output_file = os.path.join(output_dir, f"{base_name}.{target_format}")
            
            self.root.after(0, lambda: self.status_var.set("正在转换..."))
            self.root.after(0, lambda: self.progress_var.set(25))
            
            # 根据文件类型选择转换方法
            if self.current_format.startswith('image/'):
                self.convert_image(source_file, output_file, target_format)
            elif self.current_format.startswith('audio/'):
                self.convert_audio(source_file, output_file, target_format)
            elif self.current_format.startswith('video/'):
                self.convert_video(source_file, output_file, target_format)
            elif self.current_format in ['application/pdf', 'text/plain', 'text/html', 'text/markdown',
                                       'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                                       'application/msword', 'application/rtf']:
                self.convert_document(source_file, output_file, target_format)
            elif self.current_format in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                                       'application/vnd.ms-excel', 'text/csv', 'text/tab-separated-values',
                                       'application/vnd.oasis.opendocument.spreadsheet']:
                self.convert_spreadsheet(source_file, output_file, target_format)
            elif self.current_format in ['application/vnd.openxmlformats-officedocument.presentationml.presentation',
                                       'application/vnd.ms-powerpoint', 'application/vnd.oasis.opendocument.presentation']:
                self.convert_presentation(source_file, output_file, target_format)
            else:
                raise Exception(f"不支持的文件类型: {self.current_format}")
                
            self.root.after(0, lambda: self.progress_var.set(100))
            self.root.after(0, lambda: self.status_var.set("转换完成"))
            self.root.after(0, lambda: messagebox.showinfo("成功", f"文件已成功转换为: {output_file}"))
            
        except Exception as e:
            self.root.after(0, lambda: self.status_var.set("转换失败"))
            self.root.after(0, lambda: messagebox.showerror("转换错误", f"转换失败: {str(e)}"))
        finally:
            self.root.after(0, lambda: setattr(self.convert_button, 'state', 'normal'))
            
    def convert_image(self, source_file, output_file, target_format):
        """转换图片格式"""
        self.root.after(0, lambda: self.progress_var.set(50))
        
        with Image.open(source_file) as img:
            # 处理透明度
            if target_format.upper() in ['JPEG', 'JPG'] and img.mode in ('RGBA', 'LA'):
                background = Image.new('RGB', img.size, (255, 255, 255))
                background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                img = background
                
            self.root.after(0, lambda: self.progress_var.set(75))
            
            # 保存转换后的图片
            if target_format.upper() == 'JPEG':
                img.save(output_file, 'JPEG', quality=95)
            else:
                img.save(output_file, target_format.upper())
                
    def convert_audio(self, source_file, output_file, target_format):
        """转换音频格式（需要ffmpeg）"""
        self.root.after(0, lambda: self.progress_var.set(50))
        
        try:
            cmd = ['ffmpeg', '-i', source_file, '-y', output_file]
            subprocess.run(cmd, check=True, capture_output=True)
            self.root.after(0, lambda: self.progress_var.set(75))
        except subprocess.CalledProcessError as e:
            raise Exception(f"音频转换失败，请确保已安装ffmpeg: {e}")
        except FileNotFoundError:
            raise Exception("未找到ffmpeg，请先安装ffmpeg")
            
    def convert_video(self, source_file, output_file, target_format):
        """转换视频格式（需要ffmpeg）"""
        self.root.after(0, lambda: self.progress_var.set(50))
        
        try:
            cmd = ['ffmpeg', '-i', source_file, '-c:v', 'libx264', '-c:a', 'aac', '-y', output_file]
            subprocess.run(cmd, check=True, capture_output=True)
            self.root.after(0, lambda: self.progress_var.set(75))
        except subprocess.CalledProcessError as e:
            raise Exception(f"视频转换失败，请确保已安装ffmpeg: {e}")
        except FileNotFoundError:
            raise Exception("未找到ffmpeg，请先安装ffmpeg")
            
    def convert_document(self, source_file, output_file, target_format):
        """转换文档格式"""
        self.root.after(0, lambda: self.progress_var.set(50))
        
        try:
            if target_format.lower() == 'txt':
                # 处理各种格式到TXT的转换
                if self.current_format == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                    # DOCX to TXT
                    doc = Document(source_file)
                    content = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
                    with open(output_file, 'w', encoding='utf-8') as f:
                        f.write(content)
                elif self.current_format == 'application/msword':
                    # DOC to TXT (使用pypandoc)
                    content = pypandoc.convert_file(source_file, 'plain')
                    with open(output_file, 'w', encoding='utf-8') as f:
                        f.write(content)
                else:
                    # 其他格式的简单文本转换
                    with open(source_file, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()
                    with open(output_file, 'w', encoding='utf-8') as f:
                        f.write(content)
                        
            elif target_format.lower() == 'docx':
                # 转换为DOCX格式
                if self.current_format == 'text/plain':
                    # TXT to DOCX
                    doc = Document()
                    with open(source_file, 'r', encoding='utf-8') as f:
                        content = f.read()
                    doc.add_paragraph(content)
                    doc.save(output_file)
                elif self.current_format == 'application/msword':
                    # DOC to DOCX (使用pypandoc)
                    pypandoc.convert_file(source_file, 'docx', outputfile=output_file)
                else:
                    raise Exception(f"不支持从 {self.current_format} 转换为 DOCX")
                    
            elif target_format.lower() == 'html':
                # 转换为HTML格式
                if self.current_format == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                    # DOCX to HTML (使用pypandoc)
                    pypandoc.convert_file(source_file, 'html', outputfile=output_file)
                elif self.current_format == 'text/plain':
                    # TXT to HTML
                    with open(source_file, 'r', encoding='utf-8') as f:
                        content = f.read()
                    html_content = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>转换文档</title>
</head>
<body>
    <pre>{content}</pre>
</body>
</html>"""
                    with open(output_file, 'w', encoding='utf-8') as f:
                        f.write(html_content)
                else:
                    pypandoc.convert_file(source_file, 'html', outputfile=output_file)
                    
            elif target_format.lower() == 'md':
                # 转换为Markdown格式
                if self.current_format == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                    pypandoc.convert_file(source_file, 'markdown', outputfile=output_file)
                elif self.current_format == 'text/plain':
                    # TXT to MD (简单包装)
                    with open(source_file, 'r', encoding='utf-8') as f:
                        content = f.read()
                    with open(output_file, 'w', encoding='utf-8') as f:
                        f.write(f"# 转换文档\n\n{content}")
                else:
                    pypandoc.convert_file(source_file, 'markdown', outputfile=output_file)
                    
            elif target_format.lower() == 'pdf':
                # 转换为PDF格式 (使用pypandoc)
                pypandoc.convert_file(source_file, 'pdf', outputfile=output_file)
                
            else:
                raise Exception(f"不支持的文档格式: {target_format}")
                
        except Exception as e:
            if "pypandoc" in str(e).lower():
                raise Exception("需要安装pandoc来支持此转换。请访问 https://pandoc.org/installing.html")
            else:
                raise e
            
        self.root.after(0, lambda: self.progress_var.set(75))
        
    def convert_spreadsheet(self, source_file, output_file, target_format):
        """转换电子表格格式"""
        self.root.after(0, lambda: self.progress_var.set(50))
        
        try:
            # 根据源文件格式读取数据
            if self.current_format in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'application/vnd.ms-excel']:
                # Excel文件
                df = pd.read_excel(source_file)
            elif self.current_format == 'text/csv':
                # CSV文件
                df = pd.read_csv(source_file)
            elif self.current_format == 'text/tab-separated-values':
                # TSV文件
                df = pd.read_csv(source_file, sep='\t')
            else:
                raise Exception(f"不支持的源格式: {self.current_format}")
                
            self.root.after(0, lambda: self.progress_var.set(60))
            
            # 根据目标格式保存数据
            target_format = target_format.lower()
            if target_format == 'xlsx':
                df.to_excel(output_file, index=False)
            elif target_format == 'csv':
                df.to_csv(output_file, index=False, encoding='utf-8-sig')
            elif target_format == 'tsv':
                df.to_csv(output_file, sep='\t', index=False, encoding='utf-8-sig')
            elif target_format == 'xls':
                # 注意：openpyxl不直接支持.xls，这里转为xlsx
                output_file_xlsx = output_file.replace('.xls', '.xlsx')
                df.to_excel(output_file_xlsx, index=False)
                # 如果需要真正的.xls格式，需要使用xlwt库
                try:
                    import xlwt
                    workbook = xlwt.Workbook()
                    worksheet = workbook.add_sheet('Sheet1')
                    for i, col in enumerate(df.columns):
                        worksheet.write(0, i, col)
                    for i, row in df.iterrows():
                        for j, value in enumerate(row):
                            worksheet.write(i+1, j, str(value))
                    workbook.save(output_file)
                except ImportError:
                    # 如果没有xlwt，保存为xlsx格式
                    df.to_excel(output_file.replace('.xls', '.xlsx'), index=False)
                    output_file = output_file.replace('.xls', '.xlsx')
            else:
                raise Exception(f"不支持的目标格式: {target_format}")
                
        except Exception as e:
            raise Exception(f"电子表格转换失败: {str(e)}")
            
        self.root.after(0, lambda: self.progress_var.set(75))
        
    def convert_presentation(self, source_file, output_file, target_format):
        """转换演示文稿格式"""
        self.root.after(0, lambda: self.progress_var.set(50))
        
        try:
            target_format = target_format.lower()
            
            if target_format == 'pdf':
                # 转换为PDF (使用pypandoc或其他方法)
                if self.current_format == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                    # PPTX转PDF需要专门的处理
                    # 这里简化处理，提取文本内容
                    prs = Presentation(source_file)
                    content = []
                    for slide in prs.slides:
                        slide_text = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                slide_text.append(shape.text)
                        content.append('\n'.join(slide_text))
                    
                    # 创建一个简单的HTML然后转PDF
                    html_content = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>演示文稿</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 40px; }}
        .slide {{ page-break-after: always; margin-bottom: 40px; }}
        h1 {{ color: #333; }}
    </style>
</head>
<body>
"""
                    for i, slide_content in enumerate(content):
                        html_content += f'<div class="slide"><h1>幻灯片 {i+1}</h1><p>{slide_content.replace(chr(10), "<br>")}</p></div>'
                    html_content += "</body></html>"
                    
                    # 保存临时HTML文件然后转PDF
                    temp_html = output_file.replace('.pdf', '_temp.html')
                    with open(temp_html, 'w', encoding='utf-8') as f:
                        f.write(html_content)
                    
                    try:
                        pypandoc.convert_file(temp_html, 'pdf', outputfile=output_file)
                        os.remove(temp_html)  # 删除临时文件
                    except:
                        # 如果pypandoc失败，保留HTML文件
                        os.rename(temp_html, output_file.replace('.pdf', '.html'))
                        raise Exception("PDF转换失败，已保存为HTML格式")
                        
                else:
                    raise Exception("仅支持PPTX格式转换为PDF")
                    
            elif target_format == 'pptx':
                # 其他格式转PPTX的基础实现
                if self.current_format == 'application/vnd.ms-powerpoint':
                    # PPT转PPTX需要Office或LibreOffice
                    raise Exception("PPT转PPTX需要安装Microsoft Office或LibreOffice")
                else:
                    raise Exception(f"不支持从 {self.current_format} 转换为PPTX")
                    
            else:
                raise Exception(f"不支持的演示文稿格式: {target_format}")
                
        except Exception as e:
            if "pypandoc" in str(e).lower():
                raise Exception("需要安装pandoc来支持此转换。请访问 https://pandoc.org/installing.html")
            else:
                raise e
            
        self.root.after(0, lambda: self.progress_var.set(75))

def main():
    root = tk.Tk()
    app = FileConverter(root)
    root.mainloop()

if __name__ == "__main__":
    main()