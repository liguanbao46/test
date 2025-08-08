#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Office格式转换器演示脚本

演示如何使用file_converter.py中的Office格式转换功能
包括Word、Excel、PowerPoint文档的格式转换示例
"""

import os
import tempfile
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
import pandas as pd

def create_sample_word_doc(filename):
    """创建示例Word文档"""
    doc = Document()
    doc.add_heading('示例Word文档', 0)
    
    doc.add_heading('第一章', level=1)
    doc.add_paragraph('这是一个示例段落，用于演示Word文档的格式转换功能。')
    doc.add_paragraph('支持的转换格式包括：TXT, HTML, MD, PDF等。')
    
    doc.add_heading('第二章', level=1)
    p = doc.add_paragraph('这是一个')
    p.add_run('加粗文本').bold = True
    p.add_run('的示例。')
    
    # 添加表格
    table = doc.add_table(rows=1, cols=3)
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = '姓名'
    hdr_cells[1].text = '年龄'
    hdr_cells[2].text = '职业'
    
    row_cells = table.add_row().cells
    row_cells[0].text = '张三'
    row_cells[1].text = '25'
    row_cells[2].text = '程序员'
    
    doc.save(filename)
    print(f"创建示例Word文档: {filename}")

def create_sample_excel_file(filename):
    """创建示例Excel文件"""
    # 创建工作簿和工作表
    wb = Workbook()
    ws = wb.active
    ws.title = "示例数据"
    
    # 添加标题行
    headers = ['姓名', '年龄', '城市', '薪资']
    for col, header in enumerate(headers, 1):
        ws.cell(row=1, column=col, value=header)
    
    # 添加数据
    data = [
        ['张三', 25, '北京', 8000],
        ['李四', 30, '上海', 12000],
        ['王五', 28, '广州', 10000],
        ['赵六', 35, '深圳', 15000]
    ]
    
    for row_idx, row_data in enumerate(data, 2):
        for col_idx, value in enumerate(row_data, 1):
            ws.cell(row=row_idx, column=col_idx, value=value)
    
    wb.save(filename)
    print(f"创建示例Excel文件: {filename}")

def create_sample_powerpoint(filename):
    """创建示例PowerPoint文件"""
    prs = Presentation()
    
    # 幻灯片1：标题幻灯片
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Office格式转换器"
    subtitle.text = "Python实现的多格式文档转换工具"
    
    # 幻灯片2：内容幻灯片
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = '支持的格式'
    
    tf = body_shape.text_frame
    tf.text = 'Word文档格式'
    
    p = tf.add_paragraph()
    p.text = 'Excel电子表格'
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = 'PowerPoint演示文稿'
    p.level = 0
    
    # 幻灯片3：总结
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = '功能特点'
    tf = body_shape.text_frame
    tf.text = '智能格式检测'
    
    p = tf.add_paragraph()
    p.text = '高质量转换'
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = '友好用户界面'
    p.level = 0
    
    prs.save(filename)
    print(f"创建示例PowerPoint文件: {filename}")

def demo_csv_operations():
    """演示CSV文件操作"""
    # 创建示例CSV数据
    data = {
        '产品名称': ['笔记本电脑', '台式机', '平板电脑', '智能手机'],
        '价格': [5999, 3999, 2999, 1999],
        '库存': [50, 30, 80, 120],
        '分类': ['电脑', '电脑', '电脑', '手机']
    }
    
    df = pd.DataFrame(data)
    
    # 保存为不同格式
    temp_dir = tempfile.gettempdir()
    
    csv_file = os.path.join(temp_dir, 'demo_products.csv')
    excel_file = os.path.join(temp_dir, 'demo_products.xlsx')
    tsv_file = os.path.join(temp_dir, 'demo_products.tsv')
    
    # 保存为CSV
    df.to_csv(csv_file, index=False, encoding='utf-8-sig')
    print(f"创建CSV文件: {csv_file}")
    
    # 保存为Excel
    df.to_excel(excel_file, index=False)
    print(f"创建Excel文件: {excel_file}")
    
    # 保存为TSV
    df.to_csv(tsv_file, sep='\t', index=False, encoding='utf-8-sig')
    print(f"创建TSV文件: {tsv_file}")
    
    return csv_file, excel_file, tsv_file

def main():
    """主演示函数"""
    print("=" * 60)
    print("Office格式转换器演示")
    print("=" * 60)
    
    # 创建临时目录用于演示
    temp_dir = tempfile.gettempdir()
    demo_dir = os.path.join(temp_dir, 'office_converter_demo')
    
    if not os.path.exists(demo_dir):
        os.makedirs(demo_dir)
    
    print(f"演示文件将保存在: {demo_dir}")
    print()
    
    # 1. 创建Word文档示例
    print("1. 创建Word文档示例")
    word_file = os.path.join(demo_dir, 'demo_document.docx')
    create_sample_word_doc(word_file)
    print()
    
    # 2. 创建Excel文件示例
    print("2. 创建Excel文件示例")
    excel_file = os.path.join(demo_dir, 'demo_spreadsheet.xlsx')
    create_sample_excel_file(excel_file)
    print()
    
    # 3. 创建PowerPoint文件示例
    print("3. 创建PowerPoint文件示例")
    ppt_file = os.path.join(demo_dir, 'demo_presentation.pptx')
    create_sample_powerpoint(ppt_file)
    print()
    
    # 4. 创建CSV/TSV示例
    print("4. 创建电子表格格式示例")
    csv_file, excel_file2, tsv_file = demo_csv_operations()
    print()
    
    print("=" * 60)
    print("演示文件创建完成！")
    print("=" * 60)
    print("现在可以使用file_converter.py来转换这些文件：")
    print()
    print("Word文档转换示例：")
    print(f"  - 源文件: {word_file}")
    print("  - 可转换为: TXT, HTML, MD, PDF")
    print()
    print("Excel文件转换示例：")
    print(f"  - 源文件: {excel_file}")
    print("  - 可转换为: CSV, TSV, XLS")
    print()
    print("PowerPoint转换示例：")
    print(f"  - 源文件: {ppt_file}")
    print("  - 可转换为: PDF")
    print()
    print("CSV文件转换示例：")
    print(f"  - 源文件: {csv_file}")
    print("  - 可转换为: XLSX, TSV")
    print()
    print("使用方法：")
    print("1. 运行 python file_converter.py")
    print("2. 选择上述任一示例文件")
    print("3. 选择目标格式进行转换")
    print()
    print("注意：某些转换功能需要安装Pandoc")
    print("Pandoc下载地址：https://pandoc.org/installing.html")

if __name__ == "__main__":
    try:
        main()
    except ImportError as e:
        print(f"缺少必要的库: {e}")
        print("请运行以下命令安装依赖:")
        print("pip install python-docx openpyxl python-pptx pandas")
    except Exception as e:
        print(f"演示脚本执行出错: {e}")